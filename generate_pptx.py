import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Tokens
    NAVY = RGBColor(11, 45, 82)
    DARK_NAVY = RGBColor(6, 28, 54)
    TEAL = RGBColor(0, 168, 150)
    GOLD = RGBColor(224, 169, 109)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(30, 41, 59)
    LIGHT_GRAY = RGBColor(241, 245, 249)
    BORDER_COLOR = RGBColor(226, 232, 240)
    LIGHT_TEAL = RGBColor(230, 247, 245)

    def add_header(slide, title_text, category_text="영광군의회 열린소통 ON 구축 제안"):
        # Header background bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()

        # Gold accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.333), Inches(0.08))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = GOLD
        stripe.line.fill.background()

        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(10), Inches(0.3))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = GOLD

        # Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(11.5), Inches(0.55))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_NAVY
    bg1.line.fill.background()

    # Decorative shape
    dec = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(7.5), Inches(0), Inches(5.833), Inches(7.5))
    dec.fill.solid()
    dec.fill.fore_color.rgb = NAVY
    dec.line.fill.background()
    dec.rotation = 180

    # Badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(4.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(224, 169, 109)
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "🏛️ 제9대·10대 영광군의회 맞춤 소통 제안"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = DARK_NAVY
    p_b.alignment = PP_ALIGN.CENTER

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.0), Inches(2.2))
    tf_t = t_box.text_frame
    p1 = tf_t.paragraphs[0]
    p1.text = "영광군의회 군민소통 전용 웹플랫폼"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p2 = tf_t.add_paragraph()
    p2.text = "'열린소통 ON' 구축 및 도입 제안서"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.space_before = Pt(12)

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(10.0), Inches(0.8))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "\"군민의 소리를 담다, 함께 만드는 영광의 미래\""
    p_sub.font.size = Pt(20)
    p_sub.font.italic = True
    p_sub.font.color.rgb = TEAL

    # Presenter info
    info_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(8.0), Inches(0.6))
    tf_info = info_box.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "제안자: 글로컬소프트 | 영광군의회 실무자 제출용"
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = RGBColor(148, 163, 184)

    # ==========================================
    # SLIDE 2: Background & Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. 추진 배경 및 도입 필요성", "BACKGROUND & NEED")

    # Left Box: Problem
    box_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = WHITE
    box_left.line.color.rgb = BORDER_COLOR

    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "❌ 기존 지자체/의정 홈페이지의 한계점"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(231, 29, 54)

    items_l = [
      "단방향 정보 전달 중심: 공지사항 및 텍스트 회의록 열람 위주로 군민 참여 저조",
      "고령층 접근성 부족: 작은 글씨와 복잡한 메뉴 구조로 어르신 군민 소통 소외",
      "지역 밀착 소통 미흡: 영광 11개 읍·면별 맞춤 이슈 파악 및 의견 수렴 한계",
      "민원 처리 불투명성: 청원 접수 후 처리 단계와 안건 반영 과정 확인 어려움"
    ]
    for it in items_l:
        p = tf_l.add_paragraph()
        p.text = "• " + it
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(14)

    # Right Box: Solution
    box_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = LIGHT_TEAL
    box_right.line.color.rgb = TEAL

    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "⭕ '열린소통 ON'만의 혁신적 해결책"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL

    items_r = [
      "쌍방향 참여 포털: 군민청원 100명 동의 시 의회 안건 자동 반영 프로세스",
      "어르신 쉬운 모드 지원: 글자 120% 확대 & 고대비 모드 원클릭 전환 기능",
      "11개 읍·면 소통지도: 법성, 백수, 염산 등 동네별 맞춤 민원 시각화",
      "실시간 처리 시각화: [접수 ➔ 동의 ➔ 답변 ➔ 반영] 4단계 프로그레스 바 공개"
    ]
    for it in items_r:
        p = tf_r.add_paragraph()
        p.text = "✔ " + it
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_before = Pt(14)

    # ==========================================
    # SLIDE 3: 5 Core Differentiators
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. '열린소통 ON' 5대 핵심 차별화 기능", "CORE FEATURES")

    cards_data = [
      ("① 군민청원 100명 동의", "청원 제출 후 100명 동의 시 의장단 및 상임위 공식 안건 지정 & 실시간 진행바 제공"),
      ("② 어르신 쉬운 모드", "영광군 고령층 인구 특성을 고려한 원클릭 글자 확대 & 고대비 접근성 특화 UI"),
      ("③ 11개 읍·면 소통지도", "영광읍부터 낙월면 도서 지역까지 11개 지역별 민원 및 의정활동 핀포인트 안내"),
      ("④ 조례안 라이브 투표", "발의 중인 영광군 조례안에 군민이 직접 찬반 투표 및 파이차트 실시간 시각화"),
      ("⑤ 생중계 & AI 요약", "본회의 실시간 방송시청과 주요 발언 4초 간격 AI 라이브 요약 & 참관 예약 통합")
    ]

    left_margin = Inches(0.8)
    card_width = Inches(2.2)
    card_gap = Inches(0.2)

    for i, (ctitle, cdesc) in enumerate(cards_data):
        x = left_margin + i * (card_width + card_gap)
        c_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), card_width, Inches(4.8))
        c_shape.fill.solid()
        c_shape.fill.fore_color.rgb = WHITE
        c_shape.line.color.rgb = TEAL if i == 0 else BORDER_COLOR

        tf = c_shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY
        p2.space_before = Pt(14)

    # ==========================================
    # SLIDE 4: Comparison Table
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. 기존 홈페이지 vs '열린소통 ON' 비교", "COMPARISON MATRIX")

    table_shape = slide4.shapes.add_table(6, 3, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.7)
    table.columns[2].width = Inches(4.833)

    headers = ["구분", "기존 영광군의회 홈페이지 ❌", "새로 구축한 '열린소통 ON' ⭕"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if j != 2 else TEAL
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    rows_data = [
      ("소통 방식", "단방향 정보 전달 (공지사항, 회의록 단순 조회)", "쌍방향 참여형 (군민청원 동의, 조례안 찬반 직접 투표)"),
      ("어르신 접근성", "작은 글씨, 복잡한 메뉴 (어르신 이용 어려움)", "어르신 쉬운 모드 (글자 120% 확대, 고대비) 원클릭 지원"),
      ("지역 밀착도", "영광군 전체 통합 게시판 위주", "영광 11개 읍·면 소통지도 (동네별 이슈 핀포인트)"),
      ("처리 투명성", "민원 제출 후 처리 과정 파악 불가", "[접수➔동의➔답변➔반영] 4단계 실시간 진행바 공개"),
      ("회의록 요약", "난해한 전문 회의록 검색 방식", "AI 실시간 발언 요약 & 1:1 의원 직접 건의함 제공")
    ]

    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if j == 0:
                p.font.bold = True
                p.font.color.rgb = NAVY
                p.alignment = PP_ALIGN.CENTER
            elif j == 2:
                p.font.bold = True
                p.font.color.rgb = NAVY
            else:
                p.font.color.rgb = DARK_GRAY

    # ==========================================
    # SLIDE 5: Active Council Roster Integration
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. 제10대 영광군의회 현역 의장단 & 의원 맞춤 연동", "COUNCIL ROSTER")

    info_t = slide5.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5))
    tf_it = info_t.text_frame
    p_it = tf_it.paragraphs[0]
    p_it.text = "※ 영광군의회 공식 홈페이지(https://www.yeonggwang.go.kr/ygcouncil) 제10대 최신 의원 정보 반영 완료"
    p_it.font.size = Pt(13)
    p_it.font.bold = True
    p_it.font.color.rgb = TEAL

    members_data = [
      ("조일영 의장", "가선거구 (영광·대마·묘량·불갑)", "제10대 전반기 의장"),
      ("김관필 부의장", "나선거구 (백수·홍농·법성)", "제10대 전반기 부의장"),
      ("정용호 의원", "가선거구", "의회운영위원장"),
      ("김홍재 의원", "가선거구", "자치행정위원장"),
      ("박진구 의원", "나선거구", "산업건설위원장"),
      ("임영민 의원", "가선거구", "영광군의회 의원"),
      ("정선우 의원", "나선거구", "영광군의회 의원"),
      ("김선옥 의원", "비례대표", "영광군의회 의원")
    ]

    for idx, (mname, mdist, mrole) in enumerate(members_data):
        row = idx // 4
        col = idx % 4
        x = Inches(0.8) + col * Inches(2.95)
        y = Inches(2.1) + row * Inches(2.4)

        m_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.8), Inches(2.1))
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = WHITE
        m_box.line.color.rgb = BORDER_COLOR

        tf = m_box.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = mrole
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = TEAL

        p2 = tf.add_paragraph()
        p2.text = mname
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = mdist
        p3.font.size = Pt(11)
        p3.font.color.rgb = DARK_GRAY
        p3.space_before = Pt(4)

        p4 = tf.add_paragraph()
        p4.text = "✉️ 1:1 직접 건의함 탑재"
        p4.font.size = Pt(11)
        p4.font.bold = True
        p4.font.color.rgb = GOLD
        p4.space_before = Pt(8)

    # ==========================================
    # SLIDE 6: Admin Portal System
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. 의회 사무과 전용 통합 관리자 시스템 (admin.html)", "ADMIN SYSTEM")

    admin_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    admin_box.fill.solid()
    admin_box.fill.fore_color.rgb = WHITE
    admin_box.line.color.rgb = BORDER_COLOR

    tf_a = admin_box.text_frame
    tf_a.word_wrap = True

    p = tf_a.paragraphs[0]
    p.text = "⚙️ 실무자를 위한 원스톱 관리 대시보드 탑재"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY

    feats = [
      ("📢 청원 및 안건 관리", "접수된 군민 청원 조회, 의회 공식 답변 등록 및 100명 달성 안건 반영 처리"),
      ("✉️ 1:1 의원 건의 알림", "의원실별 접수된 건의 답변 작성 및 군민 신청자에게 SMS 알림 발송 연동"),
      ("📺 생중계 & AI 요약", "본회의 생중계 상태(LIVE/종료) 제어 및 AI 실시간 발언 요약 라이브 등록"),
      ("🗺️ 11개 읍·면 이슈 갱신", "각 읍면별 소통 지도 내 최신 민원 및 현장 의정소식 1초 만에 업데이트")
    ]

    for title, desc in feats:
        p = tf_a.add_paragraph()
        p.text = "• " + title + " : " + desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(16)

    # ==========================================
    # SLIDE 7: Mobile App & Expected Effects
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. 모바일 앱 확장성 및 구축 기대효과", "EXPECTED IMPACT")

    # Left: App Extension
    app_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    app_card.fill.solid()
    app_card.fill.fore_color.rgb = WHITE
    app_card.line.color.rgb = BORDER_COLOR

    tf_app = app_card.text_frame
    tf_app.word_wrap = True
    p = tf_app.paragraphs[0]
    p.text = "📱 모바일 앱(App) 출시 확장성"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEAL

    app_items = [
      "PWA (Progressive Web App): 앱스토어 설치 없이 스마트폰 홈 화면에 '영광군의회 앱' 1초 추가",
      "iOS & Android 정식 앱 출시: 하이브리드(Capacitor) 변환으로 구글/애플 스토어 정식 등록 가능",
      "스마트폰 푸시 알림 (Push): 청원 답변 등록 알림, 본회의 생중계 시작 알림 자동 발송"
    ]
    for it in app_items:
        p = tf_app.add_paragraph()
        p.text = "✔ " + it
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(16)

    # Right: Impact
    imp_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    imp_card.fill.solid()
    imp_card.fill.fore_color.rgb = DARK_NAVY
    imp_card.line.fill.background()

    tf_imp = imp_card.text_frame
    tf_imp.word_wrap = True
    p = tf_imp.paragraphs[0]
    p.text = "📈 영광군의회 3대 기대 효과"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD

    imp_items = [
      "1. 전남 1등 소통 의회 브랜딩: 군민 중심의 열린 소통 포털 운영 성과 홍보",
      "2. 전 계층 포용적 의정: 청년층(온라인 청원/투표) + 어르신(쉬운 모드) 모두 포용",
      "3. 의정 활동 성과 극대화: 1:1 열린 의원실 및 읍면 지도를 통한 공약 이행률 홍보"
    ]
    for it in imp_items:
        p = tf_imp.add_paragraph()
        p.text = it
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.space_before = Pt(16)

    # ==========================================
    # SLIDE 8: Closing & Demo Info
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = DARK_NAVY
    bg8.line.fill.background()

    close_box = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf_close = close_box.text_frame
    
    p = tf_close.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = tf_close.add_paragraph()
    p2.text = "영광군의회 군민소통 전용 웹플랫폼 '열린소통 ON'과 함께\n군민 중심의 참여 의정을 만드시길 기원합니다."
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(20)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf_close.add_paragraph()
    p3.text = "🌐 데모 시연: http://localhost:8085/ (관리자: /admin.html)\n📂 깃허브 저장소: https://github.com/shokun2002-maker/ygcouncil\n제작/지원: 글로컬소프트"
    p3.font.size = Pt(15)
    p3.font.color.rgb = TEAL
    p3.space_before = Pt(30)
    p3.alignment = PP_ALIGN.CENTER

    output_path = "/Users/glocalsoft/Desktop/코딩/ygcouncil/yeonggwang_proposal.pptx"
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
